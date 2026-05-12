"""Build Hackathon_Final_Presentation.pptx styled after portfolio/index.html.

Design system (from portfolio/index.html):
  - Surface #0A0A0A, Surface alt #111111
  - Accent #FF5500
  - Text #FFFFFF, Muted #999999
  - Display: Bebas Neue, Body: DM Sans
  - 16px rounded cards on subtle white fills with hairline borders
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --- Design tokens ---
SURFACE   = RGBColor(0x0A, 0x0A, 0x0A)
SURFACE_2 = RGBColor(0x11, 0x11, 0x11)
CARD_BG   = RGBColor(0x16, 0x16, 0x16)
ACCENT    = RGBColor(0xFF, 0x55, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x99, 0x99, 0x99)
BORDER    = RGBColor(0x26, 0x26, 0x26)

DISPLAY = "Bebas Neue"
BODY    = "DM Sans"

IMG_DIR = Path("/home/kaarthikeya/Hackathon-main/PDF/Paper/images")
OUT     = Path("/home/kaarthikeya/Hackathon-main/PDF/PPT/Hackathon_Final_Presentation.pptx")

# --- Presentation: 16:9 widescreen (13.333 x 7.5 in) ---
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = SURFACE
    bg.shadow.inherit = False
    return s


def set_text(tf, text, *, font=BODY, size=14, color=WHITE, bold=False,
             italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             tracking=0, line_spacing=1.15):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.text = ""
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if tracking:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(tracking))


def add_text_box(slide, x, y, w, h, text, **kw):
    tb = slide.shapes.add_textbox(x, y, w, h)
    set_text(tb.text_frame, text, **kw)
    return tb


def add_rect(slide, x, y, w, h, *, fill=CARD_BG, border=BORDER,
             border_w=0.75, radius=True):
    shape_t = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_t, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = border
        sh.line.width = Pt(border_w)
    sh.shadow.inherit = False
    if radius:
        # tighten the corner to ~ 8-12 px feel rather than chubby default
        sh.adjustments[0] = 0.08
    return sh


def add_pill(slide, x, y, w, h, text, *, fill=None, color=WHITE, size=10):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.5
    if fill is None:
        sh.fill.solid(); sh.fill.fore_color.rgb = SURFACE_2
        sh.line.color.rgb = BORDER
        sh.line.width = Pt(0.75)
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    set_text(tf, text, font=BODY, size=size, color=color, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=200)
    return sh


def add_brand(slide, x, y, w, h, *, size=20):
    """Display-font brand with smaller muted (SuraksaMarga.ai) suffix."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.0
    r1 = p.add_run()
    r1.text = "SurakṣāMārga.ai"
    r1.font.name = DISPLAY
    r1.font.size = Pt(size)
    r1.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = "  (SuraksaMarga.ai)"
    r2.font.name = BODY
    r2.font.size = Pt(max(8, int(size * 0.45)))
    r2.font.color.rgb = MUTED
    r2.font.bold = True
    return tb


def add_eyebrow(slide, x, y, text):
    add_text_box(slide, x, y, Inches(6), Inches(0.3), text,
                 font=BODY, size=9, color=MUTED, bold=True, tracking=300)


def add_section_index(slide, x, y, idx_text):
    """Big orange #0X label in display font."""
    add_text_box(slide, x, y, Inches(2), Inches(0.7), idx_text,
                 font=DISPLAY, size=36, color=ACCENT, line_spacing=1.0)


def add_divider(slide, x, y, w, *, color=BORDER):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(0.75)


def add_navbar(slide, current="01"):
    """Mimic top navbar: brand left, page index right."""
    add_brand(slide, Inches(0.6), Inches(0.35), Inches(6), Inches(0.4),
              size=20)
    add_text_box(slide, SW - Inches(2.6), Inches(0.4), Inches(2), Inches(0.3),
                 f"5G INNOVATION HACKATHON  ·  {current}", font=BODY, size=9,
                 color=MUTED, bold=True, tracking=300, align=PP_ALIGN.RIGHT)
    add_divider(slide, Inches(0.6), Inches(0.85), SW - Inches(1.2))


def add_image_card(slide, x, y, w, h, img_path, *, pad=Inches(0.18)):
    add_rect(slide, x, y, w, h, fill=SURFACE_2, border=BORDER)
    if not Path(img_path).exists():
        return
    inner_w = w - 2 * pad
    inner_h = h - 2 * pad
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    ar_img = iw / ih
    ar_box = inner_w / inner_h
    if ar_img > ar_box:
        nw = inner_w
        nh = int(nw / ar_img)
    else:
        nh = inner_h
        nw = int(nh * ar_img)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    slide.shapes.add_picture(str(img_path), nx, ny, width=nw, height=nh)


def add_footer(slide, page):
    add_divider(slide, Inches(0.6), SH - Inches(0.55), SW - Inches(1.2))
    add_text_box(slide, Inches(0.6), SH - Inches(0.45), Inches(6), Inches(0.3),
                 "SURAKṢĀMĀRGA.AI  ·  WOMEN-FIRST SAFE NAVIGATION",
                 font=BODY, size=8, color=MUTED, bold=True, tracking=300)
    add_text_box(slide, SW - Inches(2.6), SH - Inches(0.45), Inches(2),
                 Inches(0.3), f"{page} / 09",
                 font=BODY, size=8, color=MUTED, bold=True, tracking=300,
                 align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = add_slide()

# subtle accent hairline at top
add_rect(s, 0, 0, SW, Inches(0.06), fill=ACCENT, border=None, radius=False)

# brand
add_brand(s, Inches(0.6), Inches(0.4), Inches(7), Inches(0.4), size=22)
add_text_box(s, SW - Inches(3.5), Inches(0.45), Inches(3), Inches(0.3),
             "5G INNOVATION HACKATHON 2026", font=BODY, size=9,
             color=MUTED, bold=True, tracking=300, align=PP_ALIGN.RIGHT)

# eyebrow
add_text_box(s, Inches(0.6), Inches(1.7), Inches(8), Inches(0.3),
             "A 5G-ENABLED AI SYSTEM FOR WOMEN SAFETY NAVIGATION",
             font=BODY, size=10, color=MUTED, bold=True, tracking=300)

# Massive heading (left column only)
add_text_box(s, Inches(0.55), Inches(2.05), Inches(8), Inches(1.4),
             "SAFETY", font=DISPLAY, size=96, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.55), Inches(3.15), Inches(8), Inches(1.4),
             "OVER", font=DISPLAY, size=96, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.55), Inches(4.25), Inches(8), Inches(1.4),
             "SPEED.", font=DISPLAY, size=96, color=ACCENT,
             line_spacing=0.92)

# tagline right
add_text_box(s, Inches(8.7), Inches(2.2), Inches(4.2), Inches(0.6),
             "Great navigation should feel protective.",
             font=BODY, size=14, color=WHITE, bold=True, line_spacing=1.25)
add_text_box(s, Inches(8.7), Inches(2.95), Inches(4.2), Inches(1.6),
             "We re-engineer urban routing around the realities women face, "
             "weighing 157,160 real crime records to navigate around danger, "
             "not just toward a destination.",
             font=BODY, size=10, color=MUTED, line_spacing=1.45)

# divider
add_divider(s, Inches(0.6), Inches(5.95), SW - Inches(1.2),
            color=RGBColor(0x33, 0x33, 0x33))

# service strip
strip_y = Inches(6.15)
strip_items = [
    ("01", "Safe Route Engine"),
    ("02", "Crime-Weighted Risk"),
    ("03", "5G SOS Layer"),
    ("04", "Live Safety Loop"),
]
col_w = (SW - Inches(1.2)) / 4
for i, (idx, label) in enumerate(strip_items):
    cx = Inches(0.6) + col_w * i
    add_text_box(s, cx, strip_y, Inches(2), Inches(0.5), idx,
                 font=DISPLAY, size=28, color=ACCENT, line_spacing=1.0)
    add_text_box(s, cx, strip_y + Inches(0.55), Inches(2.6), Inches(0.4),
                 label, font=BODY, size=11, color=WHITE, bold=True)

# byline
add_text_box(s, Inches(0.6), SH - Inches(0.5), Inches(8), Inches(0.3),
             "KAARTHIKEYA PAHALWAN  ·  NISARGA S",
             font=BODY, size=9, color=MUTED, bold=True, tracking=300)


# =====================================================================
# SLIDE 2 — INTRODUCTION
# =====================================================================
s = add_slide()
add_navbar(s, "02")
add_section_index(s, Inches(0.6), Inches(1.05), "01")
add_eyebrow(s, Inches(1.5), Inches(1.32), "INTRODUCTION")

add_text_box(s, Inches(0.6), Inches(1.75), Inches(7.5), Inches(1.3),
             "NAVIGATION", font=DISPLAY, size=58, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(2.45), Inches(7.5), Inches(1.3),
             "WAS NEVER BUILT", font=DISPLAY, size=58, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(3.15), Inches(7.5), Inches(1.3),
             "FOR HER.", font=DISPLAY, size=58, color=ACCENT,
             line_spacing=0.92)

intro_body = (
    "Modern mapping apps optimise for journey time, exposing women to "
    "isolated shortcuts, desolate backstreets and unlit alleys, especially "
    "after dark.\n\n"
    "SuraksaMarga.ai bridges geospatial routing with real-world danger "
    "prevention. It actively protects women by prioritising localised "
    "safety intelligence over arbitrary travel speed."
)
add_text_box(s, Inches(0.6), Inches(4.4), Inches(7), Inches(2.3),
             intro_body, font=BODY, size=12, color=MUTED, line_spacing=1.55)

# image card
add_image_card(s, Inches(8.1), Inches(1.75), Inches(4.7), Inches(4.95),
               IMG_DIR / "SurakṣāMārga.ai landing page.png")

add_footer(s, "02")


# =====================================================================
# SLIDE 3 — PROBLEM STATEMENT
# =====================================================================
s = add_slide()
add_navbar(s, "03")
add_section_index(s, Inches(0.6), Inches(1.05), "02")
add_eyebrow(s, Inches(1.5), Inches(1.32), "PROBLEM STATEMENT")

add_text_box(s, Inches(0.6), Inches(1.75), Inches(12), Inches(1.2),
             "THE MAP IGNORES", font=DISPLAY, size=52, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.2),
             "WHAT WOMEN ALREADY KNOW.", font=DISPLAY, size=52, color=ACCENT,
             line_spacing=0.92)

# three problem cards
cards = [
    ("Time-Only Models",
     "Commercial systems apply rigid pedestrian models focused entirely on "
     "reducing time and distance, never on safety."),
    ("Forced Unsafe Shortcuts",
     "Women are routed into inherently threatening terrain without any "
     "advanced warning of crime hotspots or isolation."),
    ("No Protective Layer",
     "Users must exit the map to formulate emergency communications when "
     "danger is imminent, a critical systemic failure."),
]
card_y = Inches(4.0)
card_h = Inches(2.6)
gap = Inches(0.25)
card_w = (SW - Inches(1.2) - gap * 2) / 3
for i, (title, body) in enumerate(cards):
    cx = Inches(0.6) + (card_w + gap) * i
    add_rect(s, cx, card_y, card_w, card_h)
    add_text_box(s, cx + Inches(0.35), card_y + Inches(0.3),
                 Inches(0.6), Inches(0.4), f"0{i+1}",
                 font=DISPLAY, size=22, color=ACCENT)
    add_text_box(s, cx + Inches(0.35), card_y + Inches(0.75),
                 card_w - Inches(0.7), Inches(0.5), title,
                 font=BODY, size=15, color=WHITE, bold=True)
    add_text_box(s, cx + Inches(0.35), card_y + Inches(1.2),
                 card_w - Inches(0.7), card_h - Inches(1.4), body,
                 font=BODY, size=10.5, color=MUTED, line_spacing=1.5)

add_footer(s, "03")


# =====================================================================
# SLIDE 4 — OBJECTIVES
# =====================================================================
s = add_slide()
add_navbar(s, "04")
add_section_index(s, Inches(0.6), Inches(1.05), "03")
add_eyebrow(s, Inches(1.5), Inches(1.32), "OBJECTIVES")

add_text_box(s, Inches(0.6), Inches(1.75), Inches(12), Inches(1.4),
             "RAPID. SECURE.", font=DISPLAY, size=62, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(2.55), Inches(12), Inches(1.4),
             "WOMEN-FIRST TRANSIT.", font=DISPLAY, size=62, color=ACCENT,
             line_spacing=0.92)

objectives = [
    ("01", "Safe Route Recommendation",
     "Map the safest traversal, explicitly avoiding identified danger zones."),
    ("02", "Women-Centric Crime Weighting",
     "Penalise routes featuring histories of gender-based crime incidents."),
    ("03", "SOS Emergency System",
     "Trigger live location tracking and responder alerts seamlessly."),
    ("04", "5G Low-Latency Capability",
     "Utilise 5G for immediate safety-system distress broadcasts."),
]
card_y = Inches(4.0)
card_h = Inches(2.6)
gap = Inches(0.22)
card_w = (SW - Inches(1.2) - gap * 3) / 4
for i, (idx, title, body) in enumerate(objectives):
    cx = Inches(0.6) + (card_w + gap) * i
    add_rect(s, cx, card_y, card_w, card_h)
    add_text_box(s, cx + Inches(0.3), card_y + Inches(0.3),
                 Inches(1.2), Inches(0.6), f"{idx}",
                 font=DISPLAY, size=24, color=ACCENT)
    add_text_box(s, cx + Inches(0.3), card_y + Inches(0.85),
                 card_w - Inches(0.6), Inches(0.85), title,
                 font=BODY, size=13, color=WHITE, bold=True, line_spacing=1.2)
    add_text_box(s, cx + Inches(0.3), card_y + Inches(1.7),
                 card_w - Inches(0.6), card_h - Inches(1.9), body,
                 font=BODY, size=10, color=MUTED, line_spacing=1.5)

add_footer(s, "04")


# =====================================================================
# SLIDE 5 — EXISTING WORK / RESEARCH GAP
# =====================================================================
s = add_slide()
add_navbar(s, "05")
add_section_index(s, Inches(0.6), Inches(1.05), "04")
add_eyebrow(s, Inches(1.5), Inches(1.32), "EXISTING WORK & RESEARCH GAP")

add_text_box(s, Inches(0.6), Inches(1.75), Inches(7.5), Inches(1.1),
             "FAST, NOT SAFE.", font=DISPLAY, size=46, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(2.4), Inches(7.5), Inches(1.1),
             "WE BRIDGE THE GAP.", font=DISPLAY, size=46, color=ACCENT,
             line_spacing=0.92)

gap_body = (
    "Google Maps and OSRM rely exclusively on the \"fastest route\" "
    "mentality. They possess no capacity to distinguish a secure, well-lit "
    "main avenue from an inherently unsafe, desolate alleyway.\n\n"
    "Gender-specific risk is totally absent from modern mapping tools. "
    "Existing solutions treat all pedestrians uniformly. SuraksaMarga.ai "
    "directly enforces safety constraints first and speed constraints second."
)
add_text_box(s, Inches(0.6), Inches(3.7), Inches(7.2), Inches(3.2),
             gap_body, font=BODY, size=12, color=MUTED, line_spacing=1.55)

# right comparison block
cmp_x = Inches(8.3); cmp_w = Inches(4.5); cmp_y = Inches(1.75)
add_rect(s, cmp_x, cmp_y, cmp_w, Inches(5.0))
# Two-column header aligned to the value columns below
inner_x = cmp_x + Inches(0.35)
inner_w = cmp_w - Inches(0.7)
col_w = (inner_w - Inches(0.15)) / 2
left_x = inner_x
right_x = inner_x + col_w + Inches(0.15)

add_text_box(s, left_x, cmp_y + Inches(0.3), col_w, Inches(0.32),
             "EXISTING", font=BODY, size=10, color=MUTED,
             bold=True, tracking=300, align=PP_ALIGN.LEFT)
add_text_box(s, right_x, cmp_y + Inches(0.3), col_w, Inches(0.32),
             "OURS", font=BODY, size=10, color=ACCENT,
             bold=True, tracking=300, align=PP_ALIGN.LEFT)
add_divider(s, inner_x, cmp_y + Inches(0.7), inner_w, color=BORDER)

rows = [
    ("Routing basis", "Time / distance", "Crime-weighted safety"),
    ("Gender awareness", "None", "Women-first weighting"),
    ("Night-time logic", "Static", "Active multipliers"),
    ("SOS integration", "External app", "Built into route"),
    ("Latency target", "Best-effort", "5G edge ~10 ms"),
]
ry = cmp_y + Inches(0.85)
for label, old, new in rows:
    add_text_box(s, inner_x, ry, inner_w, Inches(0.28), label.upper(),
                 font=BODY, size=8.5, color=MUTED, bold=True, tracking=200)
    add_text_box(s, left_x, ry + Inches(0.28), col_w, Inches(0.32),
                 old, font=BODY, size=11, color=MUTED, align=PP_ALIGN.LEFT)
    add_text_box(s, right_x, ry + Inches(0.28), col_w, Inches(0.32),
                 new, font=BODY, size=11, color=ACCENT, bold=True,
                 align=PP_ALIGN.LEFT)
    ry += Inches(0.78)
    add_divider(s, inner_x, ry - Inches(0.05), inner_w)

add_footer(s, "05")


# =====================================================================
# SLIDE 6 — PROPOSED METHODOLOGY
# =====================================================================
s = add_slide()
add_navbar(s, "06")
add_section_index(s, Inches(0.6), Inches(1.05), "05")
add_eyebrow(s, Inches(1.5), Inches(1.32), "PROPOSED METHODOLOGY")

add_text_box(s, Inches(0.6), Inches(1.75), Inches(12), Inches(1.2),
             "SIX SUBSYSTEMS. ONE SAFE JOURNEY.",
             font=DISPLAY, size=44, color=WHITE, line_spacing=0.95)

methods = [
    ("01", "Routing Engine",
     "Prioritises safest paths over shortest distance by penalising "
     "high-risk segments derived from crime data."),
    ("02", "Risk Scoring",
     "Crimes targeting women receive maximum analytical weight, guaranteed "
     "circumvention of historically hostile streets."),
    ("03", "Temporal Intelligence",
     "Night-time multipliers redirect predictive paths heavily toward "
     "well-lit, trafficked thoroughfares."),
    ("04", "SOS Emergency",
     "Long-press bypasses navigation to initiate live geo-tracking and "
     "transmit high-priority alerts to responders."),
    ("05", "5G Safety Layer",
     "Removes data-transmission limits, split-second distress beacons "
     "without buffer wait times."),
    ("06", "Scalable Scope",
     "Demonstrated in Bengaluru on 157,160 records; architecture designed "
     "to expand to additional cities."),
]

# 3 columns x 2 rows of small cards on left, big image on right
grid_x = Inches(0.6); grid_y = Inches(2.95)
gw = Inches(7.8); gh = Inches(3.85)
cols, rows_n = 3, 2
gap = Inches(0.18)
cw = (gw - gap * (cols - 1)) / cols
ch = (gh - gap * (rows_n - 1)) / rows_n
for i, (idx, title, body) in enumerate(methods):
    r = i // cols; c = i % cols
    x = grid_x + (cw + gap) * c
    y = grid_y + (ch + gap) * r
    add_rect(s, x, y, cw, ch)
    add_text_box(s, x + Inches(0.22), y + Inches(0.18),
                 Inches(0.9), Inches(0.4), f"{idx}",
                 font=DISPLAY, size=18, color=ACCENT)
    add_text_box(s, x + Inches(0.22), y + Inches(0.6),
                 cw - Inches(0.44), Inches(0.4), title,
                 font=BODY, size=11.5, color=WHITE, bold=True)
    add_text_box(s, x + Inches(0.22), y + Inches(1.0),
                 cw - Inches(0.44), ch - Inches(1.1), body,
                 font=BODY, size=8.5, color=MUTED, line_spacing=1.45)

# Big image on right
add_image_card(s, Inches(8.7), Inches(2.95), Inches(4.1), Inches(3.85),
               IMG_DIR / "SurakṣāMārga.ai SAFETY LOOP CLOSURE Flow:START → Route Selected → Tracking → Arrival → SAFE CONFIRMATION.png")
add_text_box(s, Inches(8.7), Inches(2.6), Inches(4.1), Inches(0.3),
             "SAFETY LOOP CLOSURE  ·  END-TO-END FLOW",
             font=BODY, size=8.5, color=MUTED, bold=True, tracking=250)

add_footer(s, "06")


# =====================================================================
# SLIDE 7 — EXPERIMENTS & RESULTS
# =====================================================================
s = add_slide()
add_navbar(s, "07")
add_section_index(s, Inches(0.6), Inches(1.05), "06")
add_eyebrow(s, Inches(1.5), Inches(1.32), "EXPERIMENTS & RESULTS")

add_text_box(s, Inches(0.6), Inches(1.75), Inches(7.6), Inches(1.1),
             "5G EDGE CHANGES", font=DISPLAY, size=42, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(2.4), Inches(7.6), Inches(1.1),
             "EVERYTHING.", font=DISPLAY, size=42, color=ACCENT,
             line_spacing=0.92)

# table card
tx = Inches(0.6); ty = Inches(3.65); tw = Inches(7.6); th = Inches(3.2)
add_rect(s, tx, ty, tw, th)
add_text_box(s, tx + Inches(0.35), ty + Inches(0.25),
             tw - Inches(0.7), Inches(0.3),
             "EMERGENCY RESPONSE LATENCY",
             font=BODY, size=9.5, color=MUTED, bold=True, tracking=300)

# header row
hdr_y = ty + Inches(0.7)
col_xs = [tx + Inches(0.35), tx + Inches(2.8), tx + Inches(4.4),
          tx + Inches(6.0)]
col_ws = [Inches(2.4), Inches(1.55), Inches(1.55), Inches(1.55)]
headers = ["Scenario", "3G", "4G", "5G Edge"]
for i, h in enumerate(headers):
    align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
    color = MUTED if i < 3 else ACCENT
    add_text_box(s, col_xs[i], hdr_y, col_ws[i], Inches(0.35), h,
                 font=BODY, size=11, color=color, bold=True, align=align,
                 tracking=150)
add_divider(s, tx + Inches(0.35), hdr_y + Inches(0.42), tw - Inches(0.7),
            color=BORDER)

table_rows = [
    ("Route Safety Scoring", "~700 ms", "~250 ms", "~50 ms"),
    ("SOS Distress Broadcast", "~400 ms", "~150 ms", "~10 ms"),
    ("Live Tracker Updates", "Every 30 s", "Every 15 s", "Every 2 s"),
]
ry = hdr_y + Inches(0.6)
for row in table_rows:
    for i, val in enumerate(row):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        color = WHITE if i == 0 else (ACCENT if i == 3 else MUTED)
        bold = i == 0 or i == 3
        add_text_box(s, col_xs[i], ry, col_ws[i], Inches(0.4), val,
                     font=BODY, size=11.5, color=color, bold=bold,
                     align=align)
    ry += Inches(0.45)
    add_divider(s, tx + Inches(0.35), ry, tw - Inches(0.7),
                color=RGBColor(0x1c, 0x1c, 0x1c))

# takeaway
add_text_box(s, tx + Inches(0.35), ty + th - Inches(0.55),
             tw - Inches(0.7), Inches(0.4),
             "5G bypasses 4G connectivity ceilings, directly accelerating "
             "emergency response for women in transit.",
             font=BODY, size=9.5, color=MUTED, italic=True, line_spacing=1.4)

# image right
add_image_card(s, Inches(8.4), Inches(3.65), Inches(4.4), Inches(3.2),
               IMG_DIR / "SurakṣāMārga.ai Complete safe route map.png")
add_text_box(s, Inches(8.4), Inches(3.3), Inches(4.4), Inches(0.3),
             "COMPLETE ROUTE MATRIX  ·  PROTECTIVE INTERVENTIONS",
             font=BODY, size=8.5, color=MUTED, bold=True, tracking=250)

add_footer(s, "07")


# =====================================================================
# SLIDE 8 — NOVELTY
# =====================================================================
s = add_slide()
add_navbar(s, "08")
add_section_index(s, Inches(0.6), Inches(1.05), "07")
add_eyebrow(s, Inches(1.5), Inches(1.32), "NOVELTY")

add_text_box(s, Inches(0.6), Inches(1.75), Inches(7.5), Inches(1.1),
             "BUILT UNIQUELY", font=DISPLAY, size=46, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(2.4), Inches(7.5), Inches(1.1),
             "FOR PROTECTION.", font=DISPLAY, size=46, color=ACCENT,
             line_spacing=0.92)

novelties = [
    ("First", "Women-first safety routing system."),
    ("Priority", "Safety prioritised over speed."),
    ("Integrated", "SOS + navigation in a single layer."),
    ("Real-Time", "5G-enabled real-time protection."),
]
ny = Inches(3.6); nh = Inches(0.78)
gap = Inches(0.14)
row_w = Inches(7.7)
for i, (k, v) in enumerate(novelties):
    y = ny + (nh + gap) * i
    add_rect(s, Inches(0.6), y, row_w, nh)
    add_text_box(s, Inches(0.85), y + Inches(0.18), Inches(2.0),
                 Inches(0.5), k.upper(),
                 font=DISPLAY, size=20, color=ACCENT)
    add_text_box(s, Inches(3.0), y + Inches(0.22), row_w - Inches(2.6),
                 Inches(0.5), v,
                 font=BODY, size=12, color=WHITE, bold=True,
                 line_spacing=1.3)

# image right
add_image_card(s, Inches(8.7), Inches(1.75), Inches(4.1), Inches(5.05),
               IMG_DIR / "SurakṣāMārga.ai route suggestion card.png")

add_footer(s, "08")


# =====================================================================
# SLIDE 9 — CONCLUSION
# =====================================================================
s = add_slide()
add_navbar(s, "09")
add_section_index(s, Inches(0.6), Inches(1.05), "08")
add_eyebrow(s, Inches(1.5), Inches(1.32), "CONCLUSION")

add_text_box(s, Inches(0.6), Inches(1.75), Inches(12), Inches(1.2),
             "NAVIGATION,", font=DISPLAY, size=54, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(2.45), Inches(12), Inches(1.2),
             "REBUILT AS SAFETY", font=DISPLAY, size=54, color=WHITE,
             line_spacing=0.92)
add_text_box(s, Inches(0.6), Inches(3.15), Inches(12), Inches(1.2),
             "INFRASTRUCTURE.", font=DISPLAY, size=54, color=ACCENT,
             line_spacing=0.92)

# left body
body_text = (
    "By shifting the foundational routing algorithm to prioritise "
    "protection over proximity, SuraksaMarga.ai solves a critical "
    "real-world safety problem, transforming navigation from a "
    "convenience tool into real-time safety infrastructure for women."
)
add_text_box(s, Inches(0.6), Inches(4.95), Inches(7.5), Inches(1.6),
             body_text, font=BODY, size=12, color=MUTED, line_spacing=1.55)

# future enhancement card
fx = Inches(8.4); fy = Inches(4.7); fw = Inches(4.4); fh = Inches(2.05)
add_rect(s, fx, fy, fw, fh)
add_text_box(s, fx + Inches(0.3), fy + Inches(0.25), Inches(2),
             Inches(0.3), "FUTURE ENHANCEMENT",
             font=BODY, size=9.5, color=ACCENT, bold=True, tracking=300)
add_text_box(s, fx + Inches(0.3), fy + Inches(0.6),
             fw - Inches(0.6), fh - Inches(0.75),
             "Real-time data integration enables predictive heat-mapping, "
             "anticipating risk zones dynamically by combining historical "
             "patterns with live situational awareness.",
             font=BODY, size=10.5, color=WHITE, line_spacing=1.5)

# closing strip — orange accent line
add_rect(s, 0, SH - Inches(0.06), SW, Inches(0.06),
         fill=ACCENT, border=None, radius=False)
add_text_box(s, Inches(0.6), SH - Inches(0.45), Inches(8), Inches(0.3),
             "THANK YOU  ·  KAARTHIKEYA PAHALWAN  ·  NISARGA S",
             font=BODY, size=9, color=MUTED, bold=True, tracking=300)
add_text_box(s, SW - Inches(2.6), SH - Inches(0.45), Inches(2),
             Inches(0.3), "09 / 09",
             font=BODY, size=9, color=MUTED, bold=True, tracking=300,
             align=PP_ALIGN.RIGHT)


# --- Save ---
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Wrote: {OUT}")
print(f"Slides: {len(prs.slides)}")
